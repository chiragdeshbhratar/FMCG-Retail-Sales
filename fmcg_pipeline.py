"""
FMCG Retail Sales Data Pipeline & Validation Framework
=====================================================

Author: [Your Name]
Date: July 2025
Description: End-to-end data pipeline for FMCG retail sales data validation and reporting

Key Features:
- Automated data ingestion and simulation
- Comprehensive data validation with business rules
- SQL database integration with error handling
- Automated reporting (Excel, PowerPoint)
- Email alerting system for data quality issues
- Performance monitoring and logging

Tech Stack: Python, SQLite, pandas, openpyxl, python-pptx, smtplib
"""

import pandas as pd
import numpy as np
import sqlite3
import logging
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, timedelta
import os
import json
from typing import Dict, List, Tuple, Optional
import warnings
warnings.filterwarnings('ignore')

# Third-party libraries for reporting
try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, Fill, PatternFill, Border, Side
    from openpyxl.chart import BarChart, Reference
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(f"Warning: Some reporting libraries not installed: {e}")
    print("Install with: pip install openpyxl python-pptx")

class FMCGDataPipeline:
    """
    Main pipeline class for FMCG retail sales data processing
    """
    
    def __init__(self, config_path: str = "config.json"):
        """Initialize pipeline with configuration"""
        self.config = self._load_config(config_path)
        self.db_path = self.config.get('database', {}).get('path', 'fmcg_sales.db')
        self.setup_logging()
        self.validation_results = {}
        
    def _load_config(self, config_path: str) -> Dict:
        """Load configuration from JSON file"""
        default_config = {
            "database": {
                "path": "fmcg_sales.db",
                "table_name": "sales_data"
            },
            "validation": {
                "price_change_threshold": 0.5,
                "volume_outlier_threshold": 3,
                "required_fields": ["sku", "date", "volume", "price", "region"]
            },
            "reporting": {
                "output_dir": "reports",
                "excel_file": "sales_validation_report.xlsx",
                "ppt_file": "sales_summary.pptx"
            },
            "email": {
                "smtp_server": "smtp.gmail.com",
                "smtp_port": 587,
                "sender_email": "your_email@company.com",
                "recipients": ["manager@company.com"]
            }
        }
        
        if os.path.exists(config_path):
            with open(config_path, 'r') as f:
                user_config = json.load(f)
                default_config.update(user_config)
        
        return default_config
    
    def setup_logging(self):
        """Setup logging configuration"""
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(levelname)s - %(message)s',
            handlers=[
                logging.FileHandler('fmcg_pipeline.log'),
                logging.StreamHandler()
            ]
        )
        self.logger = logging.getLogger(__name__)
    
    def simulate_sales_data(self, num_records: int = 10000) -> pd.DataFrame:
        """
        Simulate realistic FMCG retail sales data
        """
        self.logger.info(f"Generating {num_records} sales records...")
        
        # Product categories and SKUs
        categories = ['Beverages', 'Snacks', 'Dairy', 'Personal Care', 'Household']
        brands = ['BrandA', 'BrandB', 'BrandC', 'BrandD', 'BrandE']
        
        # Generate SKUs
        skus = []
        for cat in categories:
            for brand in brands:
                for i in range(1, 6):  # 5 variants per brand-category
                    skus.append(f"{cat[:3].upper()}-{brand}-{i:03d}")
        
        # Regions
        regions = ['North', 'South', 'East', 'West', 'Central']
        
        # Generate date range (last 2 years)
        end_date = datetime.now()
        start_date = end_date - timedelta(days=730)
        date_range = pd.date_range(start=start_date, end=end_date, freq='D')
        
        # Generate synthetic data
        np.random.seed(42)  # For reproducibility
        
        data = []
        for i in range(num_records):
            sku = np.random.choice(skus)
            category = sku.split('-')[0]
            brand = sku.split('-')[1]
            
            # Base price varies by category
            base_prices = {
                'BEV': 2.50, 'SNA': 1.75, 'DAI': 3.25, 
                'PER': 4.50, 'HOU': 5.75
            }
            base_price = base_prices.get(category, 3.00)
            
            # Add some randomness to price
            price = base_price * np.random.uniform(0.8, 1.4)
            
            # Volume depends on price and promotions
            is_promo = np.random.choice([True, False], p=[0.2, 0.8])
            promo_multiplier = np.random.uniform(1.5, 2.5) if is_promo else 1.0
            
            base_volume = np.random.poisson(50)
            volume = int(base_volume * promo_multiplier)
            
            # Promotional price reduction
            if is_promo:
                price *= np.random.uniform(0.7, 0.9)
            
            record = {
                'sku': sku,
                'category': category,
                'brand': brand,
                'date': np.random.choice(date_range),
                'region': np.random.choice(regions),
                'volume': volume,
                'price': round(price, 2),
                'promotion': is_promo,
                'revenue': round(volume * price, 2),
                'distribution': np.random.uniform(0.1, 1.0)  # % of stores carrying product
            }
            
            data.append(record)
        
        # Introduce some data quality issues for testing
        df = pd.DataFrame(data)
        
        # Add duplicates (2% of records)
        duplicate_count = int(len(df) * 0.02)
        duplicate_indices = np.random.choice(df.index, duplicate_count, replace=False)
        duplicates = df.loc[duplicate_indices].copy()
        df = pd.concat([df, duplicates], ignore_index=True)
        
        # Add missing values (1% of records)
        missing_count = int(len(df) * 0.01)
        missing_indices = np.random.choice(df.index, missing_count, replace=False)
        df.loc[missing_indices, 'price'] = np.nan
        
        # Add outliers (extreme price changes)
        outlier_count = int(len(df) * 0.005)
        outlier_indices = np.random.choice(df.index, outlier_count, replace=False)
        df.loc[outlier_indices, 'price'] *= np.random.uniform(5, 10)
        
        self.logger.info(f"Generated {len(df)} records with quality issues included")
        return df
    
    def validate_data(self, df: pd.DataFrame) -> Dict:
        """
        Comprehensive data validation with business rules
        """
        self.logger.info("Starting data validation...")
        
        validation_results = {
            'total_records': len(df),
            'issues': {
                'duplicates': [],
                'missing_values': [],
                'price_outliers': [],
                'volume_outliers': [],
                'invalid_dates': [],
                'negative_values': []
            },
            'summary': {}
        }
        
        # 1. Check for duplicates
        duplicate_mask = df.duplicated(subset=['sku', 'date', 'region'], keep=False)
        duplicates = df[duplicate_mask].index.tolist()
        validation_results['issues']['duplicates'] = duplicates
        
        # 2. Check for missing values
        required_fields = self.config['validation']['required_fields']
        for field in required_fields:
            if field in df.columns:
                missing_mask = df[field].isna()
                missing_indices = df[missing_mask].index.tolist()
                validation_results['issues']['missing_values'].extend([
                    {'field': field, 'index': idx} for idx in missing_indices
                ])
        
        # 3. Price outliers (using IQR method)
        if 'price' in df.columns:
            Q1 = df['price'].quantile(0.25)
            Q3 = df['price'].quantile(0.75)
            IQR = Q3 - Q1
            lower_bound = Q1 - 1.5 * IQR
            upper_bound = Q3 + 1.5 * IQR
            
            price_outliers = df[(df['price'] < lower_bound) | (df['price'] > upper_bound)].index.tolist()
            validation_results['issues']['price_outliers'] = price_outliers
        
        # 4. Volume outliers
        if 'volume' in df.columns:
            volume_mean = df['volume'].mean()
            volume_std = df['volume'].std()
            threshold = self.config['validation']['volume_outlier_threshold']
            
            volume_outliers = df[
                abs(df['volume'] - volume_mean) > threshold * volume_std
            ].index.tolist()
            validation_results['issues']['volume_outliers'] = volume_outliers
        
        # 5. Invalid dates (future dates)
        if 'date' in df.columns:
            future_dates = df[df['date'] > datetime.now()].index.tolist()
            validation_results['issues']['invalid_dates'] = future_dates
        
        # 6. Negative values check
        numeric_columns = ['price', 'volume', 'revenue']
        for col in numeric_columns:
            if col in df.columns:
                negative_indices = df[df[col] < 0].index.tolist()
                validation_results['issues']['negative_values'].extend([
                    {'field': col, 'index': idx} for idx in negative_indices
                ])
        
        # Calculate summary statistics
        total_issues = sum(len(issues) if isinstance(issues, list) else 0 
                          for issues in validation_results['issues'].values())
        
        validation_results['summary'] = {
            'total_issues': total_issues,
            'clean_records': len(df) - total_issues,
            'data_quality_score': round((len(df) - total_issues) / len(df) * 100, 2),
            'duplicate_rate': round(len(duplicates) / len(df) * 100, 2),
            'missing_value_rate': round(len(validation_results['issues']['missing_values']) / len(df) * 100, 2),
            'outlier_rate': round((len(price_outliers) + len(volume_outliers)) / len(df) * 100, 2)
        }
        
        self.validation_results = validation_results
        self.logger.info(f"Validation complete. Data quality score: {validation_results['summary']['data_quality_score']}%")
        
        return validation_results
    
    def create_database(self):
        """Create SQLite database and tables"""
        self.logger.info("Setting up database...")
        
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.cursor()
            
            # Create main sales table
            cursor.execute('''
                CREATE TABLE IF NOT EXISTS sales_data (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    sku TEXT NOT NULL,
                    category TEXT,
                    brand TEXT,
                    date DATE NOT NULL,
                    region TEXT NOT NULL,
                    volume INTEGER,
                    price REAL,
                    promotion BOOLEAN,
                    revenue REAL,
                    distribution REAL,
                    created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
                )
            ''')
            
            # Create validation log table
            cursor.execute('''
                CREATE TABLE IF NOT EXISTS validation_log (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    validation_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                    total_records INTEGER,
                    issues_found INTEGER,
                    data_quality_score REAL,
                    validation_details TEXT
                )
            ''')
            
            conn.commit()
        
        self.logger.info("Database setup complete")
    
    def load_data_to_db(self, df: pd.DataFrame, clean_only: bool = False):
        """Load data to database with option to load only clean records"""
        self.logger.info(f"Loading data to database (clean_only={clean_only})...")
        
        if clean_only and hasattr(self, 'validation_results'):
            # Remove problematic records
            all_issue_indices = set()
            for issue_list in self.validation_results['issues'].values():
                if isinstance(issue_list, list):
                    for item in issue_list:
                        if isinstance(item, dict):
                            all_issue_indices.add(item['index'])
                        else:
                            all_issue_indices.add(item)
            
            clean_df = df.drop(index=list(all_issue_indices))
            self.logger.info(f"Loading {len(clean_df)} clean records out of {len(df)} total")
        else:
            clean_df = df
        
        # Load to database
        with sqlite3.connect(self.db_path) as conn:
            clean_df.to_sql('sales_data', conn, if_exists='replace', index=False)
            
            # Log validation results
            if hasattr(self, 'validation_results'):
                cursor = conn.cursor()
                cursor.execute('''
                    INSERT INTO validation_log 
                    (total_records, issues_found, data_quality_score, validation_details)
                    VALUES (?, ?, ?, ?)
                ''', (
                    self.validation_results['total_records'],
                    self.validation_results['summary']['total_issues'],
                    self.validation_results['summary']['data_quality_score'],
                    json.dumps(self.validation_results['summary'])
                ))
                conn.commit()
        
        self.logger.info("Data loaded to database successfully")
    
    def generate_excel_report(self, df: pd.DataFrame):
        """Generate comprehensive Excel report"""
        self.logger.info("Generating Excel report...")
        
        # Create output directory
        output_dir = self.config['reporting']['output_dir']
        os.makedirs(output_dir, exist_ok=True)
        
        excel_file = os.path.join(output_dir, self.config['reporting']['excel_file'])
        
        with pd.ExcelWriter(excel_file, engine='openpyxl') as writer:
            # Summary sheet
            summary_data = {
                'Metric': ['Total Records', 'Clean Records', 'Data Quality Score (%)', 
                          'Duplicate Rate (%)', 'Missing Value Rate (%)', 'Outlier Rate (%)'],
                'Value': [
                    self.validation_results['summary']['total_issues'] + self.validation_results['summary']['clean_records'],
                    self.validation_results['summary']['clean_records'],
                    self.validation_results['summary']['data_quality_score'],
                    self.validation_results['summary']['duplicate_rate'],
                    self.validation_results['summary']['missing_value_rate'],
                    self.validation_results['summary']['outlier_rate']
                ]
            }
            
            summary_df = pd.DataFrame(summary_data)
            summary_df.to_excel(writer, sheet_name='Summary', index=False)
            
            # Issues details
            issues_data = []
            for issue_type, issues in self.validation_results['issues'].items():
                if isinstance(issues, list) and issues:
                    for issue in issues:
                        if isinstance(issue, dict):
                            issues_data.append({
                                'Issue Type': issue_type,
                                'Field': issue.get('field', 'N/A'),
                                'Record Index': issue.get('index', 'N/A')
                            })
                        else:
                            issues_data.append({
                                'Issue Type': issue_type,
                                'Field': 'Multiple',
                                'Record Index': issue
                            })
            
            if issues_data:
                issues_df = pd.DataFrame(issues_data)
                issues_df.to_excel(writer, sheet_name='Issues', index=False)
            
            # Sample of clean data
            clean_indices = set(range(len(df)))
            all_issue_indices = set()
            
            for issue_list in self.validation_results['issues'].values():
                if isinstance(issue_list, list):
                    for item in issue_list:
                        if isinstance(item, dict):
                            all_issue_indices.add(item['index'])
                        else:
                            all_issue_indices.add(item)
            
            clean_indices = clean_indices - all_issue_indices
            clean_sample = df.iloc[list(clean_indices)].head(1000)
            clean_sample.to_excel(writer, sheet_name='Clean Data Sample', index=False)
            
            # Regional analysis
            regional_stats = df.groupby('region').agg({
                'volume': ['sum', 'mean'],
                'revenue': ['sum', 'mean'],
                'price': 'mean'
            }).round(2)
            
            regional_stats.columns = ['Total Volume', 'Avg Volume', 'Total Revenue', 'Avg Revenue', 'Avg Price']
            regional_stats.to_excel(writer, sheet_name='Regional Analysis')
        
        self.logger.info(f"Excel report generated: {excel_file}")
    
    def generate_powerpoint_summary(self):
        """Generate PowerPoint summary presentation"""
        self.logger.info("Generating PowerPoint summary...")
        
        output_dir = self.config['reporting']['output_dir']
        ppt_file = os.path.join(output_dir, self.config['reporting']['ppt_file'])
        
        # Create presentation
        prs = Presentation()
        
        # Slide 1: Title
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        
        title.text = "FMCG Sales Data Quality Report"
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Slide 2: Executive Summary
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Executive Summary"
        
        summary_text = f"""
        • Total Records Processed: {self.validation_results['summary']['total_issues'] + self.validation_results['summary']['clean_records']:,}
        • Data Quality Score: {self.validation_results['summary']['data_quality_score']}%
        • Clean Records: {self.validation_results['summary']['clean_records']:,}
        • Issues Identified: {self.validation_results['summary']['total_issues']:,}
        
        Key Findings:
        • Duplicate Rate: {self.validation_results['summary']['duplicate_rate']}%
        • Missing Value Rate: {self.validation_results['summary']['missing_value_rate']}%
        • Outlier Rate: {self.validation_results['summary']['outlier_rate']}%
        """
        
        content = slide2.placeholders[1]
        content.text = summary_text
        
        # Slide 3: Recommendations
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "Recommendations"
        
        recommendations = """
        Data Quality Improvements:
        • Implement real-time validation at data entry points
        • Establish automated duplicate detection processes
        • Create data quality monitoring dashboards
        • Develop SLA for data quality metrics (target: >95%)
        
        Process Enhancements:
        • Schedule daily validation runs
        • Implement exception handling for outliers
        • Create automated alerts for quality threshold breaches
        • Establish data governance protocols
        """
        
        slide3.placeholders[1].text = recommendations
        
        prs.save(ppt_file)
        self.logger.info(f"PowerPoint summary generated: {ppt_file}")
    
    def send_alert_email(self, validation_results: Dict):
        """Send email alert if data quality issues are found"""
        
        if validation_results['summary']['data_quality_score'] < 90:  # Alert threshold
            self.logger.info("Data quality below threshold - sending alert email...")
            
            try:
                # Email configuration
                smtp_config = self.config['email']
                
                # Create message
                msg = MIMEMultipart()
                msg['From'] = smtp_config['sender_email']
                msg['To'] = ', '.join(smtp_config['recipients'])
                msg['Subject'] = f"Data Quality Alert - Score: {validation_results['summary']['data_quality_score']}%"
                
                body = f"""
                FMCG Sales Data Quality Alert
                
                Data Quality Score: {validation_results['summary']['data_quality_score']}%
                Total Records: {validation_results['total_records']:,}
                Issues Found: {validation_results['summary']['total_issues']:,}
                
                Issue Breakdown:
                - Duplicates: {validation_results['summary']['duplicate_rate']}%
                - Missing Values: {validation_results['summary']['missing_value_rate']}%
                - Outliers: {validation_results['summary']['outlier_rate']}%
                
                Please review the detailed report for corrective actions.
                
                Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
                """
                
                msg.attach(MIMEText(body, 'plain'))
                
                # Note: Email sending would require actual SMTP credentials
                # For demo purposes, we'll just log the attempt
                self.logger.info("Email alert prepared (actual sending requires SMTP configuration)")
                
            except Exception as e:
                self.logger.error(f"Failed to send email alert: {str(e)}")
        else:
            self.logger.info("Data quality within acceptable range - no alert needed")
    
    def run_pipeline(self, num_records: int = 10000):
        """Run the complete data pipeline"""
        self.logger.info("Starting FMCG Data Pipeline...")
        
        try:
            # 1. Data Generation
            df = self.simulate_sales_data(num_records)
            
            # 2. Data Validation
            validation_results = self.validate_data(df)
            
            # 3. Database Operations
            self.create_database()
            self.load_data_to_db(df, clean_only=True)
            
            # 4. Generate Reports
            self.generate_excel_report(df)
            self.generate_powerpoint_summary()
            
            # 5. Send Alerts
            self.send_alert_email(validation_results)
            
            self.logger.info("Pipeline completed successfully!")
            
            return {
                'status': 'success',
                'records_processed': len(df),
                'data_quality_score': validation_results['summary']['data_quality_score'],
                'reports_generated': True
            }
            
        except Exception as e:
            self.logger.error(f"Pipeline failed: {str(e)}")
            return {
                'status': 'failed',
                'error': str(e)
            }

# Example usage and testing
if __name__ == "__main__":
    # Create and run pipeline
    pipeline = FMCGDataPipeline()
    result = pipeline.run_pipeline(num_records=5000)
    
    print("\n" + "="*50)
    print("FMCG DATA PIPELINE RESULTS")
    print("="*50)
    print(f"Status: {result['status']}")
    if result['status'] == 'success':
        print(f"Records Processed: {result['records_processed']:,}")
        print(f"Data Quality Score: {result['data_quality_score']}%")
        print(f"Reports Generated: {result['reports_generated']}")
    else:
        print(f"Error: {result['error']}")
    
    # Display sample validation results
    if hasattr(pipeline, 'validation_results'):
        print("\nValidation Summary:")
        for metric, value in pipeline.validation_results['summary'].items():
            print(f"  {metric}: {value}")